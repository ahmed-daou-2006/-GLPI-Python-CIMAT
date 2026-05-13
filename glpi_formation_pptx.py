# glpi_formation_pptx.py — Présentation Formation GLPI CIMAT
# Projet CIMAT Béni Mellal

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from datetime import datetime

# ── Couleurs CIMAT ────────────────────────────────────────
BLEU_FONCE  = RGBColor(31,  78,  121)
BLEU_MOY    = RGBColor(46, 117, 182)
BLEU_CLAIR  = RGBColor(214, 228, 240)
VERT        = RGBColor(0,  176,  80)
ORANGE      = RGBColor(197,  90,  17)
BLANC       = RGBColor(255, 255, 255)
NOIR        = RGBColor(0,    0,   0)
GRIS        = RGBColor(89,  89,  89)

def set_bg(slide, r, g, b):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)

def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False,
                color=NOIR, align=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def slide_titre(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 31, 78, 121)
    add_rect(slide, 0, 2.5, 10, 2.5, BLEU_MOY)
    add_textbox(slide, "GLPI CIMAT", 0.5, 2.6, 9, 1,
                font_size=44, bold=True,
                color=BLANC, align=PP_ALIGN.CENTER)
    add_textbox(slide,
                "Formation — Automatisation avec Python",
                0.5, 3.4, 9, 0.7,
                font_size=20, color=BLEU_CLAIR,
                align=PP_ALIGN.CENTER)
    add_textbox(slide,
                "CIMAT Béni Mellal  |  Ahmed Daou  |  " +
                datetime.now().strftime("%d/%m/%Y"),
                0.5, 6.5, 9, 0.5,
                font_size=12, color=BLEU_CLAIR,
                align=PP_ALIGN.CENTER)

def slide_plan(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 245, 248, 252)
    add_rect(slide, 0, 0, 10, 1.1, BLEU_FONCE)
    add_textbox(slide, "Plan de la formation",
                0.3, 0.15, 9, 0.8,
                font_size=28, bold=True, color=BLANC)
    items = [
        ("01", "Introduction à GLPI",            "5 min"),
        ("02", "Interface et navigation",         "5 min"),
        ("03", "Gestion des tickets",             "5 min"),
        ("04", "Scripts Python — démonstration",  "10 min"),
        ("05", "Rapports Excel et alertes",       "5 min"),
        ("06", "Résultats et bénéfices CIMAT",    "5 min"),
    ]
    for i, (num, titre, duree) in enumerate(items):
        y = 1.3 + i * 0.9
        add_rect(slide, 0.4, y, 0.7, 0.7, BLEU_MOY)
        add_textbox(slide, num, 0.4, y+0.1, 0.7, 0.5,
                    font_size=16, bold=True,
                    color=BLANC, align=PP_ALIGN.CENTER)
        add_textbox(slide, titre, 1.3, y+0.1, 6.5, 0.5,
                    font_size=15, bold=True, color=BLEU_FONCE)
        add_rect(slide, 7.8, y+0.1, 1.5, 0.5, BLEU_CLAIR)
        add_textbox(slide, duree, 7.8, y+0.15, 1.5, 0.4,
                    font_size=12, color=BLEU_FONCE,
                    align=PP_ALIGN.CENTER)

def slide_intro_glpi(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 245, 248, 252)
    add_rect(slide, 0, 0, 10, 1.1, BLEU_FONCE)
    add_textbox(slide, "01 — C'est quoi GLPI ?",
                0.3, 0.15, 9, 0.8,
                font_size=26, bold=True, color=BLANC)
    add_textbox(slide,
                "GLPI = Gestionnaire Libre de Parc Informatique",
                0.5, 1.3, 9, 0.6,
                font_size=18, bold=True, color=BLEU_FONCE)
    points = [
        ("", "Gérer tout le matériel informatique de CIMAT"),
        ("", "Créer et suivre les tickets de support"),
        ("", "Générer des rapports et statistiques"),
        ("", "Recevoir des alertes automatiques par email"),
        ("", "Automatisé avec des scripts Python"),
    ]
    for i, (icon, text) in enumerate(points):
        y = 2.1 + i * 0.85
        add_rect(slide, 0.4, y, 8.8, 0.65, BLEU_CLAIR)
        add_textbox(slide, f"{icon}  {text}",
                    0.6, y+0.08, 8.5, 0.5,
                    font_size=15, color=BLEU_FONCE)

def slide_tickets(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 245, 248, 252)
    add_rect(slide, 0, 0, 10, 1.1, BLEU_FONCE)
    add_textbox(slide, "03 — Gestion des Tickets",
                0.3, 0.15, 9, 0.8,
                font_size=26, bold=True, color=BLANC)
    statuts = [
        ("Nouveau",   RGBColor(0, 176, 80)),
        ("En cours",  BLEU_MOY),
        ("En attente",ORANGE),
        ("Résolu",    RGBColor(112, 173, 71)),
        ("Clos",      GRIS),
    ]
    add_textbox(slide, "Cycle de vie d'un ticket :",
                0.5, 1.2, 9, 0.5,
                font_size=16, bold=True, color=BLEU_FONCE)
    for i, (statut, color) in enumerate(statuts):
        x = 0.5 + i * 1.85
        add_rect(slide, x, 1.8, 1.6, 0.65, color)
        add_textbox(slide, statut, x, 1.9, 1.6, 0.5,
                    font_size=12, bold=True,
                    color=BLANC, align=PP_ALIGN.CENTER)
        if i < 4:
            add_textbox(slide, "→", x+1.6, 1.95, 0.25, 0.4,
                        font_size=16, bold=True, color=BLEU_FONCE)
    priorites = [
        ("Très basse", RGBColor(0, 176, 240)),
        ("Basse",      RGBColor(0, 176, 80)),
        ("Moyenne",    ORANGE),
        ("Haute",      RGBColor(255, 0, 0)),
        ("Majeure",    RGBColor(139, 0, 0)),
    ]
    add_textbox(slide, "Niveaux de priorité :",
                0.5, 2.8, 9, 0.5,
                font_size=16, bold=True, color=BLEU_FONCE)
    for i, (prio, color) in enumerate(priorites):
        x = 0.5 + i * 1.85
        add_rect(slide, x, 3.4, 1.6, 0.65, color)
        add_textbox(slide, prio, x, 3.5, 1.6, 0.5,
                    font_size=12, bold=True,
                    color=BLANC, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.4, 4.3, 9, 1.5, RGBColor(240, 240, 240))
    add_textbox(slide, "Script Python — Créer un ticket :",
                0.6, 4.35, 8.5, 0.4,
                font_size=12, bold=True, color=BLEU_FONCE)
    add_textbox(slide,
                'session = connect_glpi()\n'
                'create_ticket(session,\n'
                '    titre="Panne reseau atelier",\n'
                '    priorite=5)',
                0.6, 4.75, 8.5, 1.0,
                font_size=11, color=NOIR,
                font_name="Courier New")

def slide_scripts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 245, 248, 252)
    add_rect(slide, 0, 0, 10, 1.1, BLEU_FONCE)
    add_textbox(slide, "04 — Scripts Python",
                0.3, 0.15, 9, 0.8,
                font_size=26, bold=True, color=BLANC)
    scripts = [
        ("glpi_tickets.py",
         "Création automatique de tickets",
         "python glpi_tickets.py",
         BLEU_MOY),
        ("glpi_export_excel.py",
         "Export rapport Excel hebdomadaire",
         "python glpi_export_excel.py",
         RGBColor(0, 176, 80)),
        ("glpi_alertes_email.py",
         "Alertes email tickets non résolus",
         "python glpi_alertes_email.py",
         ORANGE),
        ("glpi_inventaire.py",
         "Inventaire enrichi avec données réseau",
         "python glpi_inventaire.py",
         RGBColor(112, 48, 160)),
    ]
    for i, (fichier, desc, cmd, color) in enumerate(scripts):
        y = 1.3 + i * 1.3
        add_rect(slide, 0.3, y, 0.15, 1.0, color)
        add_textbox(slide, fichier, 0.6, y+0.05, 4.5, 0.45,
                    font_size=13, bold=True, color=color)
        add_textbox(slide, desc, 0.6, y+0.5, 4.5, 0.4,
                    font_size=11, color=GRIS)
        add_rect(slide, 5.3, y+0.25, 4.3, 0.45,
                 RGBColor(240, 240, 240))
        add_textbox(slide, cmd, 5.4, y+0.3, 4.2, 0.4,
                    font_size=10, color=NOIR,
                    font_name="Courier New")

def slide_excel(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 245, 248, 252)
    add_rect(slide, 0, 0, 10, 1.1, BLEU_FONCE)
    add_textbox(slide, "05 — Rapports Excel et Alertes",
                0.3, 0.15, 9, 0.8,
                font_size=26, bold=True, color=BLANC)
    add_rect(slide, 0.3, 1.2, 4.4, 5.5, RGBColor(235, 245, 235))
    add_textbox(slide, " Rapport Excel",
                0.5, 1.3, 4, 0.5,
                font_size=16, bold=True,
                color=RGBColor(0, 128, 0))
    excel_points = [
        "Généré chaque lundi à 08h00",
        "Liste complète des tickets",
        "Couleurs par statut et priorité",
        "Statistiques résumées",
        "Fichier : GLPI_Rapport_DATE.xlsx",
    ]
    for i, p in enumerate(excel_points):
        add_textbox(slide, f"• {p}", 0.5, 1.85+i*0.75,
                    4.0, 0.5, font_size=12, color=NOIR)
    add_rect(slide, 5.3, 1.2, 4.4, 5.5, RGBColor(255, 243, 230))
    add_textbox(slide, " Alertes Email",
                5.5, 1.3, 4, 0.5,
                font_size=16, bold=True, color=ORANGE)
    email_points = [
        "Envoyé chaque jour à 09h00",
        "Tickets non résolus +24h",
        "Email HTML professionnel",
        "Lien direct vers GLPI",
        "Destinataire : responsable IT",
    ]
    for i, p in enumerate(email_points):
        add_textbox(slide, f"• {p}", 5.5, 1.85+i*0.75,
                    4.0, 0.5, font_size=12, color=NOIR)
    add_rect(slide, 0.3, 6.5, 9.4, 0.7, BLEU_CLAIR)
    add_textbox(slide,
                "Planification automatique : Windows Task Scheduler",
                0.5, 6.6, 9, 0.5,
                font_size=13, bold=True, color=BLEU_FONCE,
                align=PP_ALIGN.CENTER)

def slide_benefices(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 245, 248, 252)
    add_rect(slide, 0, 0, 10, 1.1, BLEU_FONCE)
    add_textbox(slide, "06 — Résultats et Bénéfices pour CIMAT",
                0.3, 0.15, 9, 0.8,
                font_size=24, bold=True, color=BLANC)
    add_rect(slide, 0.3, 1.2, 4.4, 5.5, RGBColor(235, 245, 235))
    add_textbox(slide, " Gain de temps",
                0.5, 1.3, 4, 0.5,
                font_size=16, bold=True,
                color=RGBColor(0, 128, 0))
    gains = [
        "Tickets créés en 2 secondes",
        "Export Excel automatique",
        "Alertes sans intervention",
        "Inventaire en 1 clic",
        "Planification 24h/24",
    ]
    for i, g in enumerate(gains):
        add_textbox(slide, f"  {g}",
                    0.5, 1.9+i*0.85, 4.0, 0.6,
                    font_size=12, color=NOIR)
    add_rect(slide, 5.3, 1.2, 4.4, 5.5, RGBColor(230, 240, 255))
    add_textbox(slide, " Bénéfices CIMAT",
                5.5, 1.3, 4, 0.5,
                font_size=16, bold=True, color=BLEU_FONCE)
    benefices = [
        "Zéro ticket perdu",
        "Suivi en temps réel",
        "Rapports professionnels",
        "Réduction des pannes",
        "Meilleure communication IT",
    ]
    for i, b in enumerate(benefices):
        add_textbox(slide, f"  {b}",
                    5.5, 1.9+i*0.85, 4.0, 0.6,
                    font_size=12, color=NOIR)

def slide_fin(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, 31, 78, 121)
    add_rect(slide, 0, 2.8, 10, 1.8, BLEU_MOY)
    add_textbox(slide, "Merci !",
                0.5, 1.5, 9, 1.2,
                font_size=48, bold=True,
                color=BLANC, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Des questions ?",
                0.5, 2.9, 9, 0.8,
                font_size=24, color=BLANC,
                align=PP_ALIGN.CENTER)
    infos = [
        "GitHub : github.com/ahmed-daou-2006/-GLPI-Python-CIMAT",
        "Email  : ahmeddaou2006@gmail.com",
        "CIMAT Béni Mellal — " + datetime.now().strftime("%d/%m/%Y"),
    ]
    for i, info in enumerate(infos):
        add_textbox(slide, info, 0.5, 4.5+i*0.45, 9, 0.4,
                    font_size=12, color=BLEU_CLAIR,
                    align=PP_ALIGN.CENTER)

# ── MAIN ──────────────────────────────────────────────────
if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Génération présentation GLPI CIMAT...")

    slide_titre(prs)
    slide_plan(prs)
    slide_intro_glpi(prs)
    slide_tickets(prs)
    slide_scripts(prs)
    slide_excel(prs)
    slide_benefices(prs)
    slide_fin(prs)

    nom = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        f"Formation_GLPI_CIMAT_{datetime.now().strftime('%Y-%m-%d')}.pptx"
    )
    prs.save(nom)
    print(f" Présentation créée : {nom}")